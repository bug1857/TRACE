from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─────────────────────────────────────────────
#  COLOR PALETTE
# ─────────────────────────────────────────────
BG       = RGBColor(0x0A, 0x0E, 0x1A)   # deep navy
BG2      = RGBColor(0x0F, 0x17, 0x2A)   # card bg
ACCENT   = RGBColor(0x00, 0xC9, 0x8D)   # neon green — TRACE primary
ACCENT2  = RGBColor(0x00, 0x8B, 0xFF)   # electric blue
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xB8, 0xC5, 0xD6)   # muted text
GOLD     = RGBColor(0xF5, 0xC5, 0x18)   # highlight / winner
DARK_CARD= RGBColor(0x13, 0x1E, 0x35)   # slightly lighter card
RED      = RGBColor(0xFF, 0x4D, 0x6D)

def rgb(r,g,b): return RGBColor(r,g,b)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # truly blank

# ─────────────────────────────────────────────
#  HELPER FUNCTIONS
# ─────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w if line_w else Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=Pt(12), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def bg(slide):
    """Fill slide background dark navy."""
    add_rect(slide, 0, 0, W, H, fill=BG)

def accent_bar(slide, y=Inches(0.12), thickness=Pt(4)):
    """Top accent bar."""
    add_rect(slide, 0, y, W, thickness, fill=ACCENT)

def slide_number(slide, n, total=10):
    add_text(slide, f"{n} / {total}", Inches(12.3), Inches(7.1), Inches(1), Inches(0.3),
             size=Pt(9), color=LIGHT, align=PP_ALIGN.RIGHT)

def tag(slide, text, l, t, color=ACCENT):
    """Small pill badge."""
    w = Inches(1.8)
    h = Inches(0.28)
    add_rect(slide, l, t, w, h, fill=color)
    add_text(slide, text.upper(), l, t, w, h,
             size=Pt(8), bold=True, color=BG, align=PP_ALIGN.CENTER)

def card(slide, l, t, w, h, title=None, body_lines=None,
         accent_color=ACCENT, title_size=Pt(11), body_size=Pt(10)):
    """Dark card with optional title and bullet body."""
    add_rect(slide, l, t, w, h, fill=DARK_CARD, line_color=accent_color, line_w=Pt(0.75))
    # top accent stripe
    add_rect(slide, l, t, w, Pt(3), fill=accent_color)
    cy = t + Inches(0.18)
    if title:
        add_text(slide, title, l + Inches(0.15), cy, w - Inches(0.3), Inches(0.25),
                 size=title_size, bold=True, color=WHITE)
        cy += Inches(0.28)
    if body_lines:
        body_text = "\n".join(body_lines)
        add_text(slide, body_text, l + Inches(0.15), cy,
                 w - Inches(0.3), h - (cy - t) - Inches(0.1),
                 size=body_size, color=LIGHT)

def stat_card(slide, l, t, w, h, number, label, sub="", color=ACCENT):
    add_rect(slide, l, t, w, h, fill=DARK_CARD, line_color=color, line_w=Pt(0.75))
    add_rect(slide, l, t, w, Pt(3), fill=color)
    add_text(slide, number, l, t + Inches(0.15), w, Inches(0.55),
             size=Pt(32), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, t + Inches(0.65), w, Inches(0.28),
             size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, sub, l, t + Inches(0.90), w, Inches(0.22),
                 size=Pt(8), color=LIGHT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 1 — TITLE / COVER
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
# Background gradient effect — two overlapping rects
add_rect(s, 0, 0, Inches(7), H, fill=rgb(0x06, 0x10, 0x20))
add_rect(s, Inches(6.5), 0, Inches(6.83), H, fill=rgb(0x03, 0x18, 0x2A))

# Diagonal accent line
line_shape = s.shapes.add_shape(1, Inches(6.3), 0, Pt(2), H)
line_shape.fill.solid(); line_shape.fill.fore_color.rgb = ACCENT
line_shape.line.fill.background()

# Big "T" watermark
add_text(s, "T", Inches(7.2), Inches(-0.5), Inches(5), Inches(8),
         size=Pt(420), bold=True, color=rgb(0x00, 0x25, 0x1A), align=PP_ALIGN.LEFT)

# TRACE Logo text
add_text(s, "TRACE.", Inches(0.6), Inches(1.6), Inches(6), Inches(1.2),
         size=Pt(72), bold=True, color=WHITE)

# Green underline
add_rect(s, Inches(0.6), Inches(2.7), Inches(3.2), Pt(4), fill=ACCENT)

add_text(s, "Carbon Intelligence Platform", Inches(0.6), Inches(2.85), Inches(8), Inches(0.45),
         size=Pt(22), bold=False, color=ACCENT)

add_text(s,
    "Process Mining × Scope 3 Carbon Accounting\nfor Enterprise Supply Chains",
    Inches(0.6), Inches(3.4), Inches(6.2), Inches(0.85),
    size=Pt(15), color=LIGHT)

# Bottom meta row
add_rect(s, 0, Inches(6.8), W, Inches(0.7), fill=rgb(0x05, 0x0C, 0x18))
add_text(s, "Manipal University Jaipur  ·  Indo-Swiss Joint Research Programme",
         Inches(0.6), Inches(6.85), Inches(8), Inches(0.4),
         size=Pt(10), color=LIGHT)
add_text(s, "hackathon submission  2026", Inches(9.5), Inches(6.85), Inches(3.5), Inches(0.4),
         size=Pt(10), color=ACCENT, align=PP_ALIGN.RIGHT)

slide_number(s, 1)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 2 — THE PROBLEM
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "THE PROBLEM", Inches(0.55), Inches(0.35), Inches(5), Inches(0.3),
         size=Pt(9), bold=True, color=ACCENT)
add_text(s, "Supply chain emissions are\na regulatory blind spot.", Inches(0.55), Inches(0.6), Inches(7), Inches(1.3),
         size=Pt(38), bold=True, color=WHITE)

# Three pain point cards
pain = [
    ("~73%", "of corporate carbon\nemissions are Scope 3\n— yet most go\nunmeasured.", ACCENT2),
    ("Manual &\nBroken", "ESG teams rely on\nspreadsheets, surveys,\nand estimates.\nNo real audit trail.", RED),
    ("Regulatory\nDeadline", "BRSR, CSRD, GHG\nProtocol mandates\nrequire verifiable\nper-activity data.", GOLD),
]
for i, (num, desc, col) in enumerate(pain):
    lx = Inches(0.55 + i * 4.2)
    add_rect(s, lx, Inches(2.1), Inches(3.9), Inches(4.5), fill=DARK_CARD,
             line_color=col, line_w=Pt(1))
    add_rect(s, lx, Inches(2.1), Inches(3.9), Pt(4), fill=col)
    add_text(s, num, lx + Inches(0.2), Inches(2.3), Inches(3.5), Inches(1.0),
             size=Pt(36), bold=True, color=col)
    add_text(s, desc, lx + Inches(0.2), Inches(3.3), Inches(3.5), Inches(2.8),
             size=Pt(14), color=LIGHT)

add_text(s, "\"We don't know where our emissions come from — or why they're failing compliance.\"",
         Inches(0.55), Inches(6.8), Inches(12), Inches(0.5),
         size=Pt(11), italic=True, color=LIGHT)
slide_number(s, 2)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 3 — OUR SOLUTION
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "THE SOLUTION", Inches(0.55), Inches(0.35), Inches(5), Inches(0.3),
         size=Pt(9), bold=True, color=ACCENT)
add_text(s, "TRACE turns raw event logs into\ncertifiable carbon intelligence.", Inches(0.55), Inches(0.6), Inches(10), Inches(1.3),
         size=Pt(34), bold=True, color=WHITE)

add_text(s, "Upload any CSV logistics log  →  TRACE reconstructs every route, computes per-activity Scope 3 emissions,\ndetects policy violations, and auto-generates audit-ready BRSR / ESG compliance reports.",
         Inches(0.55), Inches(1.85), Inches(12.3), Inches(0.7),
         size=Pt(13), color=LIGHT)

# Flow diagram
steps = [
    ("01", "OCEL Upload\n& Mapping", "Drag CSV → fuzzy\ncolumn detection\nauto-maps schema", ACCENT2),
    ("02", "Process Mining\nEngine", "Reconstructs DFG\ngraph of actual\nfreight routes", ACCENT),
    ("03", "Carbon\nComputation", "Per-activity Scope 3\nemissions with\nemission factors", ACCENT),
    ("04", "Conformance\nAudit", "Table-driven rule\nengine flags every\npolicy violation", GOLD),
    ("05", "Reporting &\nCopilot", "BRSR, ESG reports +\nlocal LLM natural\nlanguage queries", RED),
]
for i, (num, title, desc, col) in enumerate(steps):
    lx = Inches(0.3 + i * 2.58)
    add_rect(s, lx, Inches(2.7), Inches(2.4), Inches(4.2), fill=DARK_CARD,
             line_color=col, line_w=Pt(0.75))
    add_rect(s, lx, Inches(2.7), Inches(2.4), Pt(3), fill=col)
    add_text(s, num, lx + Inches(0.15), Inches(2.85), Inches(0.6), Inches(0.38),
             size=Pt(20), bold=True, color=col)
    add_text(s, title, lx + Inches(0.15), Inches(3.25), Inches(2.1), Inches(0.55),
             size=Pt(12), bold=True, color=WHITE)
    add_text(s, desc, lx + Inches(0.15), Inches(3.8), Inches(2.1), Inches(2.8),
             size=Pt(10), color=LIGHT)
    # arrow between cards
    if i < 4:
        add_text(s, "→", lx + Inches(2.38), Inches(4.2), Inches(0.25), Inches(0.35),
                 size=Pt(14), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

slide_number(s, 3)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 4 — TECHNICAL ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "TECHNICAL ARCHITECTURE", Inches(0.55), Inches(0.35), Inches(6), Inches(0.3),
         size=Pt(9), bold=True, color=ACCENT)
add_text(s, "Full-stack, production-grade,\nzero external AI dependencies.", Inches(0.55), Inches(0.6), Inches(8), Inches(1.0),
         size=Pt(30), bold=True, color=WHITE)

# Tech layer boxes
layers = [
    ("FRONTEND", "Next.js 16  ·  React 19  ·  TypeScript (strict)  ·  Tailwind v4  ·  ReactFlow  ·  Recharts", ACCENT2),
    ("BACKEND API", "Python FastAPI  ·  SQLAlchemy ORM  ·  SQLite  ·  Pydantic schemas  ·  pytest (36/36 passing)", ACCENT),
    ("COMPUTATION ENGINE", "column_mapper.py (fuzzy)  ·  process_mining.py (DFG)  ·  conformance.py (table-driven rules)\ncarbon_fitness.py (CFS = idealKg/actualKg × 100)  ·  forecasting.py  ·  brsr_report.py", GOLD),
    ("AI COPILOT", "Ollama local LLM  ·  gemma3:4b / qwen2.5 / mistral  ·  Rich context injection  ·  Privacy-preserving (no cloud)", RED),
    ("MULTI-TENANCY", "Organization → Project → Workspace → AnalysisSnapshot  ·  Cascade-delete  ·  Session-persisted state", LIGHT),
]
for i, (layer, tech, col) in enumerate(layers):
    lx = Inches(0.55)
    ly = Inches(1.75 + i * 0.97)
    add_rect(s, lx, ly, Inches(12.3), Inches(0.82), fill=DARK_CARD, line_color=col, line_w=Pt(0.75))
    add_rect(s, lx, ly, Inches(1.9), Inches(0.82), fill=col)
    add_text(s, layer, lx + Inches(0.08), ly + Inches(0.18), Inches(1.75), Inches(0.45),
             size=Pt(8), bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_text(s, tech, lx + Inches(2.1), ly + Inches(0.12), Inches(10.0), Inches(0.6),
             size=Pt(10), color=WHITE)

slide_number(s, 4)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 5 — INNOVATION
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "INNOVATION", Inches(0.55), Inches(0.35), Inches(5), Inches(0.3),
         size=Pt(9), bold=True, color=ACCENT)
add_text(s, "First platform to combine\nprocess mining + real-time Scope 3.", Inches(0.55), Inches(0.6), Inches(10), Inches(1.3),
         size=Pt(30), bold=True, color=WHITE)

innovations = [
    ("🔬", "Process-Native\nCarbon Accounting", "Emissions computed from actual event log paths — not averages or estimates. Every kg is traceable to a case ID, activity, and timestamp.", ACCENT),
    ("⚡", "Fuzzy Column\nAuto-Mapping", "Zero-config ingestion: any CSV with logistics columns is auto-detected using confidence-scored fuzzy alias matching. Supports 8 field types including water, electricity, cost.", ACCENT2),
    ("🛡️", "Table-Driven\nConformance Engine", "Rule violations detected via a generalized `CONFORMANCE_RULES` table — not hardcoded logic. Pluggable: future PNML upload replaces the table entirely.", GOLD),
    ("🤖", "Privacy-First\nAI Copilot", "Local Ollama LLM with rich structured context injection (violations, suppliers, budgets, ESG). Zero data leaves the machine. Model-switching verified via network interception.", RED),
    ("📋", "Cryptographically\nSealed Reports", "BRSR reports include SHA-256 hash of all source data fields. Tamper-evident audit trail ready for regulatory submission.", ACCENT2),
    ("🏢", "Real Multi-Tenancy", "Org → Project → Workspace → AnalysisSnapshot hierarchy with cascade-delete. No fake dropdowns — every entity is SQL-backed and real.", ACCENT),
]
for i, (icon, title, desc, col) in enumerate(innovations):
    row = i // 3
    col_i = i % 3
    lx = Inches(0.4 + col_i * 4.3)
    ly = Inches(2.0 + row * 2.45)
    add_rect(s, lx, ly, Inches(4.0), Inches(2.2), fill=DARK_CARD, line_color=col, line_w=Pt(0.75))
    add_rect(s, lx, ly, Inches(4.0), Pt(3), fill=col)
    add_text(s, icon + "  " + title, lx + Inches(0.15), ly + Inches(0.12), Inches(3.7), Inches(0.55),
             size=Pt(11), bold=True, color=WHITE)
    add_text(s, desc, lx + Inches(0.15), ly + Inches(0.68), Inches(3.7), Inches(1.4),
             size=Pt(9.5), color=LIGHT)

slide_number(s, 5)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 6 — LIVE DEMO: WHAT THE PLATFORM DOES
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "PLATFORM CAPABILITIES  —  LIVE DEMO", Inches(0.55), Inches(0.35), Inches(8), Inches(0.3),
         size=Pt(9), bold=True, color=ACCENT)
add_text(s, "10 fully-functional modules,\nzero mock data.", Inches(0.55), Inches(0.6), Inches(9), Inches(1.0),
         size=Pt(32), bold=True, color=WHITE)

modules = [
    ("Executive Dashboard", "Real-time KPIs: total carbon, CFS score, violation count", ACCENT),
    ("Process Mining Graph", "Directed graph of actual vs. intended freight routes", ACCENT2),
    ("Conformance Ledger", "Per-case violation audit with carbon delta annotations", RED),
    ("Supplier Fitness", "Per-vendor CFS scoring, at-risk flagging, corrective actions", GOLD),
    ("Carbon Budget Tracker", "Monthly burn rate + predicted breach date", ACCENT),
    ("Forecasting Engine", "ARIMA + linear trend — next-month emission predictions", ACCENT2),
    ("What-If Simulator", "Modal shift scenarios (Air→Ocean) with live recalculation", GOLD),
    ("BRSR Report Generator", "One-click SHA-256 sealed compliance report", RED),
    ("ESG Scorecard", "E/S/G pillar scoring with data-completeness disclosure", ACCENT),
    ("AI Copilot", "Natural language queries over live supply chain data", ACCENT2),
]
for i, (name, desc, col) in enumerate(modules):
    row = i // 2
    col_i = i % 2
    lx = Inches(0.4 + col_i * 6.5)
    ly = Inches(1.78 + row * 0.94)
    add_rect(s, lx, ly, Inches(6.2), Inches(0.78), fill=DARK_CARD, line_color=col, line_w=Pt(0.5))
    add_rect(s, lx, ly, Pt(4), Inches(0.78), fill=col)
    add_text(s, name, lx + Inches(0.2), ly + Inches(0.06), Inches(2.8), Inches(0.3),
             size=Pt(11), bold=True, color=WHITE)
    add_text(s, desc, lx + Inches(0.2), ly + Inches(0.38), Inches(5.8), Inches(0.35),
             size=Pt(9), color=LIGHT)

slide_number(s, 6)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 7 — THE MATH: CFS & CONFORMANCE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "TECHNICAL DEPTH", Inches(0.55), Inches(0.35), Inches(5), Inches(0.3),
         size=Pt(9), bold=True, color=ACCENT)
add_text(s, "The math behind TRACE.", Inches(0.55), Inches(0.6), Inches(8), Inches(0.7),
         size=Pt(36), bold=True, color=WHITE)

# CFS Formula box
add_rect(s, Inches(0.5), Inches(1.5), Inches(5.9), Inches(2.5), fill=DARK_CARD,
         line_color=ACCENT, line_w=Pt(1))
add_rect(s, Inches(0.5), Inches(1.5), Inches(5.9), Pt(3), fill=ACCENT)
add_text(s, "Carbon Fitness Score (CFS)", Inches(0.65), Inches(1.62), Inches(5.6), Inches(0.3),
         size=Pt(11), bold=True, color=ACCENT)
add_text(s, "CFS = min( idealCarbonKg / actualCarbonKg × 100 ,  100 )",
         Inches(0.65), Inches(2.0), Inches(5.6), Inches(0.5),
         size=Pt(14), bold=True, color=WHITE)
add_text(s,
    "idealCarbonKg  = Σ (activity_emission_factor × frequency)  for compliant path only\n"
    "actualCarbonKg = Σ all observed emission events for the case\n"
    "CFS = 100   →  fully compliant route\n"
    "CFS < 100  →  excess Scope 3 from violations",
    Inches(0.65), Inches(2.55), Inches(5.6), Inches(1.3),
    size=Pt(10), color=LIGHT)

# Conformance formula box
add_rect(s, Inches(6.9), Inches(1.5), Inches(5.9), Inches(2.5), fill=DARK_CARD,
         line_color=RED, line_w=Pt(1))
add_rect(s, Inches(6.9), Inches(1.5), Inches(5.9), Pt(3), fill=RED)
add_text(s, "Conformance Rule Engine", Inches(7.05), Inches(1.62), Inches(5.6), Inches(0.3),
         size=Pt(11), bold=True, color=RED)
add_text(s, "violation  ←  activity ∈ DISALLOWED_ACTIVITIES",
         Inches(7.05), Inches(2.0), Inches(5.6), Inches(0.5),
         size=Pt(13), bold=True, color=WHITE)
add_text(s,
    "carbonDeltaKg = actualKg  −  mandatedAlternativeKg\n"
    "severity = CRITICAL  if  carbonDeltaKg > threshold_C\n"
    "severity = HIGH       if  carbonDeltaKg > threshold_H\n"
    "Table-driven: rules defined in CONFORMANCE_RULES[]",
    Inches(7.05), Inches(2.55), Inches(5.6), Inches(1.3),
    size=Pt(10), color=LIGHT)

# Bottom — demo numbers
add_text(s, "Real demo results  (trace_demo_dataset.csv  —  800 cases, 4,232 events)",
         Inches(0.55), Inches(4.2), Inches(12), Inches(0.3),
         size=Pt(10), bold=True, color=ACCENT)

stats = [
    ("368", "Violations\nDetected", ACCENT2),
    ("73.4%", "Avg Carbon\nFitness Score", ACCENT),
    ("Air Freight\nDispatch", "Top Violation\nActivity", RED),
    ("36 / 36", "Pytest\nPassing", GOLD),
]
for i, (n, l, c) in enumerate(stats):
    lx = Inches(0.5 + i * 3.2)
    stat_card(s, lx, Inches(4.65), Inches(2.9), Inches(2.4), n, l, color=c)

slide_number(s, 7)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 8 — FEASIBILITY
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "FEASIBILITY", Inches(0.55), Inches(0.35), Inches(5), Inches(0.3),
         size=Pt(9), bold=True, color=ACCENT)
add_text(s, "Production-ready today.\nEnterprise-scalable tomorrow.", Inches(0.55), Inches(0.6), Inches(9), Inches(1.1),
         size=Pt(34), bold=True, color=WHITE)

left_items = [
    ("✅ Working demo — not a prototype", "Full end-to-end pipeline from CSV to BRSR report runs in real-time on actual logistics data."),
    ("✅ Zero infrastructure dependencies", "SQLite backend, local LLM via Ollama — deployable on any machine with no cloud account required."),
    ("✅ Tested & verified", "36/36 pytest passing. Every feature verified against real data, not mocks. No fake numbers shown alongside real ones."),
    ("✅ Standards-aligned", "Output maps directly to BRSR (India MCA), CSRD (EU), and GHG Protocol Scope 3 Category 4 reporting."),
]
right_items = [
    ("📈 Scale path — Phase 2", "PostgreSQL swap (SQLAlchemy ORM already abstract). REST API already exposes all endpoints for ERP integration."),
    ("🔗 SAP / Oracle connector", "Column mapper's fuzzy-alias engine can ingest structured exports from SAP SCM or Oracle TMS with minimal config changes."),
    ("☁️ Cloud Copilot fallback", "OLLAMA_BASE_URL env var allows pointing Copilot to any self-hosted or managed LLM instance for cloud deployments."),
    ("🌐 Multi-org SaaS", "Multi-tenancy (Org → Project → Workspace) already fully built and SQL-backed. License-gating is the only addition needed."),
]
for i, (title, desc) in enumerate(left_items):
    ly = Inches(1.88 + i * 1.28)
    add_rect(s, Inches(0.5), ly, Inches(6.0), Inches(1.1), fill=DARK_CARD, line_color=ACCENT, line_w=Pt(0.5))
    add_text(s, title, Inches(0.65), ly + Inches(0.08), Inches(5.7), Inches(0.3),
             size=Pt(10), bold=True, color=WHITE)
    add_text(s, desc, Inches(0.65), ly + Inches(0.38), Inches(5.7), Inches(0.65),
             size=Pt(9.5), color=LIGHT)

for i, (title, desc) in enumerate(right_items):
    ly = Inches(1.88 + i * 1.28)
    add_rect(s, Inches(6.9), ly, Inches(6.0), Inches(1.1), fill=DARK_CARD, line_color=ACCENT2, line_w=Pt(0.5))
    add_text(s, title, Inches(7.05), ly + Inches(0.08), Inches(5.7), Inches(0.3),
             size=Pt(10), bold=True, color=WHITE)
    add_text(s, desc, Inches(7.05), ly + Inches(0.38), Inches(5.7), Inches(0.65),
             size=Pt(9.5), color=LIGHT)

slide_number(s, 8)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 9 — EXPECTED IMPACT
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "EXPECTED IMPACT", Inches(0.55), Inches(0.35), Inches(5), Inches(0.3),
         size=Pt(9), bold=True, color=ACCENT)
add_text(s, "Every supply chain audit\nbecomes a carbon audit.", Inches(0.55), Inches(0.6), Inches(9), Inches(1.1),
         size=Pt(34), bold=True, color=WHITE)

impact_stats = [
    ("73%", "of corporate\nemissions are\nScope 3 — now\nfully traceable", ACCENT),
    ("0 hrs", "manual ESG\nspreadsheet work\nto generate\nBRSR report", ACCENT2),
    ("100%", "of violations\ntraceable to\na case ID,\nactivity & time", GOLD),
    ("∞", "scalability:\nSaaS-ready\nmulti-tenant\narchitecture", RED),
]
for i, (n, l, c) in enumerate(impact_stats):
    lx = Inches(0.5 + i * 3.2)
    stat_card(s, lx, Inches(1.9), Inches(2.9), Inches(2.4), n, l, color=c)

# Impact narrative cards
narratives = [
    ("For Enterprise Supply Chain Teams",
     "Replace 3-month manual carbon audits with a 30-second CSV upload. Every logistics event becomes a verifiable data point in your ESG ledger.", ACCENT),
    ("For Regulators & Auditors",
     "SHA-256 sealed BRSR reports with full traceability matrix (metric → engine → source table → reference field). Audit-ready from day one.", GOLD),
    ("For Sustainability Officers",
     "Natural language queries over your own supply chain data — no data leaves your machine. Ask 'Which supplier is worst for Scope 3?' and get a data-backed answer.", ACCENT2),
    ("For Researchers",
     "Built on OCEL 2.0 process mining + real GHG Protocol emission factors. Academic-grade methodology packaged in an enterprise-grade interface.", RED),
]
for i, (title, desc, col) in enumerate(narratives):
    col_i = i % 2
    row = i // 2
    lx = Inches(0.5 + col_i * 6.5)
    ly = Inches(4.55 + row * 1.35)
    add_rect(s, lx, ly, Inches(6.2), Inches(1.2), fill=DARK_CARD, line_color=col, line_w=Pt(0.75))
    add_rect(s, lx, ly, Pt(4), Inches(1.2), fill=col)
    add_text(s, title, lx + Inches(0.2), ly + Inches(0.08), Inches(5.8), Inches(0.28),
             size=Pt(10), bold=True, color=WHITE)
    add_text(s, desc, lx + Inches(0.2), ly + Inches(0.38), Inches(5.8), Inches(0.75),
             size=Pt(9.5), color=LIGHT)

slide_number(s, 9)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 10 — CLOSING / THANK YOU
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s)

# Full-bleed diagonal split
add_rect(s, 0, 0, Inches(7.5), H, fill=rgb(0x05, 0x12, 0x22))
add_rect(s, Inches(7.0), 0, Inches(6.33), H, fill=rgb(0x00, 0x18, 0x10))

# Green accent line
line2 = s.shapes.add_shape(1, Inches(7.0), 0, Pt(2), H)
line2.fill.solid(); line2.fill.fore_color.rgb = ACCENT
line2.line.fill.background()

# Left side — tagline
add_text(s, "TRACE.", Inches(0.6), Inches(1.5), Inches(6), Inches(1.2),
         size=Pt(72), bold=True, color=WHITE)
add_rect(s, Inches(0.6), Inches(2.65), Inches(2.8), Pt(4), fill=ACCENT)
add_text(s, "We don't estimate.\nWe compute.", Inches(0.6), Inches(2.78), Inches(6), Inches(1.2),
         size=Pt(26), bold=True, color=ACCENT)
add_text(s,
    "Process Mining  ×  Scope 3 Carbon Accounting\nfor Enterprise Supply Chains",
    Inches(0.6), Inches(4.1), Inches(6.2), Inches(0.9),
    size=Pt(14), color=LIGHT)

# Right side — key links and contact
add_text(s, "GitHub", Inches(7.4), Inches(1.8), Inches(5.5), Inches(0.35),
         size=Pt(10), bold=True, color=ACCENT)
add_text(s, "github.com/bug1857/TRACE", Inches(7.4), Inches(2.12), Inches(5.5), Inches(0.35),
         size=Pt(14), color=WHITE)

add_text(s, "Live Demo", Inches(7.4), Inches(2.85), Inches(5.5), Inches(0.35),
         size=Pt(10), bold=True, color=ACCENT)
add_text(s, "localhost:3000  (running now)", Inches(7.4), Inches(3.15), Inches(5.5), Inches(0.35),
         size=Pt(14), color=WHITE)

add_text(s, "Research Context", Inches(7.4), Inches(3.85), Inches(5.5), Inches(0.35),
         size=Pt(10), bold=True, color=ACCENT)
add_text(s, "Indo-Swiss Joint Research Programme\nManipul University Jaipur", Inches(7.4), Inches(4.15), Inches(5.5), Inches(0.55),
         size=Pt(12), color=LIGHT)

# Gold winner badge
add_rect(s, Inches(7.4), Inches(5.5), Inches(4.5), Inches(0.65), fill=GOLD)
add_text(s, "🏆  Built to win. Ready to scale.", Inches(7.4), Inches(5.52), Inches(4.5), Inches(0.6),
         size=Pt(14), bold=True, color=BG, align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(s, 0, Inches(6.85), W, Inches(0.65), fill=rgb(0x03, 0x09, 0x14))
add_text(s, "TRACE Carbon Intelligence Platform  ·  Hackathon 2026  ·  No mock data. No fake features. No estimates.",
         Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.4),
         size=Pt(9), color=LIGHT, align=PP_ALIGN.CENTER)

slide_number(s, 10)

# ─────────────────────────────────────────────────────────────────────────────
#  SAVE
# ─────────────────────────────────────────────────────────────────────────────
OUT = "/Users/rudrapratapsingh/Desktop/TRACE/TRACE_Winner_Pitch.pptx"
prs.save(OUT)
print(f"✅ Saved: {OUT}")
print(f"   Slides: {len(prs.slides)}")
